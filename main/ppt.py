# .\python\python-3.12.3-embed-amd64\python.exe ppt.py

import os
import warnings
from tqdm.rich import tqdm
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from rich import traceback

traceback.install(
    show_locals=True,
    extra_lines=2,
    max_frames=10
)
warnings.filterwarnings("ignore")

def add_textbox_centered(slide, text, font_name, font_size, slide_width, slide_height, x_offset, y_offset):
    # 计算文本框的大小（大致估算）
    text_length = len(text)
    text_box_width = text_length * font_size * 0.1  # 估算文本框宽度
    text_box_height = font_size * 1.2  # 估算文本框高度
    
    left = (slide_width - text_box_width) / 2 + x_offset
    top = (slide_height - text_box_height) / 2 + y_offset

    text_box = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
    text_frame = text_box.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = font_size
    p.font.name = font_name
    p.alignment = PP_ALIGN.CENTER

    # 设置文本颜色和透明度
    font_color = RGBColor(255, 255, 255)  # 白色
    p.font.color.rgb = font_color

    return text_box

def read_color_list(color_file):
    colors = []
    with open(color_file, 'r') as f:
        for line in f:
            color = line.strip()
            # 如果颜色值只有3位（如#a），则补全为6位（如#aa）
            if len(color) == 4:  # 包括开头的#
                color = '#' + color[1] * 2 + color[2] * 2 + color[3] * 2
            colors.append(color)
    return colors

def parse_hex_color(hex_str):
    """ 解析颜色字符串并创建RGBColor对象 """
    if hex_str.startswith('#') and len(hex_str) in (4, 7):
        hex_str = hex_str.lstrip('#')
        if len(hex_str) == 3:
            hex_str = ''.join([c * 2 for c in hex_str])
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        return RGBColor(r, g, b)
    else:
        raise ValueError("Invalid hex color string format")

def save_ppt(prs, folder_name, filename_prefix, file_count):
    """ 保存Presentation对象到文件夹中的文件 """
    output_filename = os.path.join(folder_name, f"{filename_prefix}_{file_count}.pptx")
    prs.save(output_filename)
    print(f"生成的PPT文件保存为：{output_filename}")

def process_unicode_file():
    # 创建一个新的Presentation对象
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9比例宽度
    prs.slide_height = Inches(7.5)   # 16:9比例高度
    BottomFont = 'RHR SC Bold'
    color_file = 'blcolor.ini'
    threshold = 1114111  # 设置分割阈值

    # 读取颜色列表
    colors = read_color_list(color_file)
    color_index = 0  # 用于循环颜色列表的索引

    # 读取unicode.txt文件
    with open('unicode.txt', 'r', encoding='utf-8') as file:
        lines = file.readlines()
        first_line = lines[0].strip()
        _, font_dir, _, _ = first_line.split(',')
        font_dir = font_dir.strip('\"')

    # 创建保存文件夹
    folder_name = font_dir.replace(' ', '_').lower()  # 使用font_dir创建文件夹名
    os.makedirs(folder_name, exist_ok=True)  # 如果文件夹不存在则创建

    # 使用tqdm显示进度条
    total_slides = sum(1 for _ in open('unicode.txt', 'r', encoding='utf-8'))  # 统计总行数（即总幻灯片数）
    progress_bar = tqdm(total=total_slides, desc='PPT创建进度', unit='slide')
    file_count = 1
    slide_count = 0

    # 提前定义颜色列表长度，避免在循环中重复计算
    num_colors = len(colors)

    for line in lines:
        try:
            # 去除换行符并拆分行内容
            line = line.strip()
            char_code, font_dir, description, change_background = line.strip().split(',')
            char_code = char_code.strip('\"')
            font_dir = font_dir.strip('\"')
            description = description.strip('\"').replace('|', '\n')
            change_background = int(change_background.strip('\"'))

            # 提取和转换Unicode字符
            code_point = int(char_code.replace('U+', ''), 16)
            unicode_char = str(chr(code_point))

        # 创建一个幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白幻灯片布局

            # 删除默认的“点击此处添加标题”文本框
            for shape in slide.shapes:
                if shape.has_text_frame and "" in shape.text_frame.text:
                    slide.shapes._spTree.remove(shape._element)

            # 添加Unicode字符到幻灯片中央，并指定x和y轴偏移量
            font_name = font_dir
            font_size = Pt(350)  # 调整字体大小
            x_offset = Inches(0)  # 示例x轴偏移量
            y_offset = Inches(-0.5)  # 示例y轴偏移量

            add_textbox_centered(slide, unicode_char, font_name, font_size, prs.slide_width, prs.slide_height, x_offset, y_offset)
    
            # 添加描述文本到幻灯片左下角，包括Unicode码点和description
            unicode_code_text = f"{char_code}\n"
            full_description = unicode_code_text + description

            # 左下角文本框位置移动
            description_box = slide.shapes.add_textbox(Inches(0.1), Inches(6.0), Inches(13.13), Inches(1.5))
            description_frame = description_box.text_frame
            description_frame.text = full_description

            # 设置文本框的自动换行
            for i, paragraph in enumerate(description_frame.paragraphs):
                paragraph.space_before = Pt(0)
                paragraph.space_after = Pt(0)
                paragraph.font.size = Pt(18)
                paragraph.font.name = BottomFont  # 使用自定义字体
                paragraph.font.color.rgb = RGBColor(255, 255, 255) if i == 0 else RGBColor(0, 0, 0)

            # 根据change_background值确定是否更改背景颜色
            if change_background == 1:
                # 更改背景颜色，并循环使用颜色列表中的颜色
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = parse_hex_color(colors[color_index])
                color_index = (color_index + 1) % num_colors
            else:
                # 不更改背景颜色，沿用当前行的颜色
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = parse_hex_color(colors[color_index])

            slide_count += 1
            progress_bar.update(1)

            # 每达到阈值时保存当前ppt，并创建新的ppt文件
            if slide_count >= threshold:
                save_ppt(prs, folder_name, 'unicode_flash_mob', file_count)
                # 创建新的Presentation对象时保持颜色索引不变
                prs = Presentation()
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                slides = [prs.slides.add_slide(prs.slide_layouts[5]) for _ in range(min(threshold, total_slides - i - 1))]
                file_count += 1
                slide_count = 0

        except Exception as e:
            print(f"Error processing line: {line}")
            print(f"Error message: {str(e)}")
            print(f"请手动删除PPT内空白页")
            continue  # 继续处理下一行

    # 处理剩余的幻灯片并保存最后一个ppt文件
    if slide_count > 0:
        save_ppt(prs, folder_name, 'unicode_flash_mob', file_count)

    progress_bar.close()

if __name__ == "__main__":
    process_unicode_file()  # 运行主程序