import os
import re
from tqdm import tqdm
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL, MSO_THEME_COLOR

def copy_background(source_slide, dest_slide):
    # 检查源幻灯片是否有背景填充
    if source_slide.background.fill.type != MSO_FILL.BACKGROUND:
        fill_type = source_slide.background.fill.type
        
        if fill_type == MSO_FILL.SOLID:
            # 复制纯色背景
            fore_color = source_slide.background.fill.fore_color
            dest_slide.background.fill.solid()
            dest_slide.background.fill.fore_color.rgb = fore_color.rgb
        
        elif fill_type == MSO_FILL.GRADIENT:
            # 复制渐变背景
            dest_slide.background.fill.gradient()
            dest_slide.background.fill.gradient_stops = source_slide.background.fill.gradient_stops
        
        elif fill_type == MSO_FILL.PICTURE:
            # 复制图片背景
            dest_slide.background.fill.background_picture()
            dest_slide.background.fill._blob = source_slide.background.fill._blob

def merge_pptx(files, output_file):
    total_slides = sum(len(Presentation(f).slides) for f in files)
    
    with tqdm(total=total_slides, desc="合并进度") as pbar:
        # 处理第一个文件
        merged_pres = Presentation(files[0])
        pbar.update(len(merged_pres.slides))
        merged_pres.save(output_file)
        
        # 处理剩余文件
        for file in files[1:]:
            # 重新打开合并后的文件
            merged_pres = Presentation(output_file)
            current_pres = Presentation(file)
            
            for slide in current_pres.slides:
                blank_slide_layout = merged_pres.slide_layouts[6]
                new_slide = merged_pres.slides.add_slide(blank_slide_layout)
                
                # 复制背景
                copy_background(slide, new_slide)
                
                # 复制形状
                for shape in slide.shapes:
                    el = shape.element
                    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
                
                pbar.update(1)
            
            # 保存当前进度
            merged_pres.save(output_file)
    
    # 返回最终的幻灯片数量
    final_pres = Presentation(output_file)
    return len(final_pres.slides)

def get_sorted_pptx_files(directory):
    # 正则表达式匹配 "unicode_flash_mob_" 开头，中间是数字，".pptx" 结尾的文件
    pattern = re.compile(r'unicode_flash_mob_(\d+)\.pptx$')
    # 获取所有匹配的文件
    matching_files = [f for f in os.listdir(directory) if pattern.match(f)]
    # 按照中间的数字排序
    sorted_files = sorted(matching_files, key=lambda f: int(pattern.match(f).group(1)))
    # 返回完整的文件路径
    return [os.path.join(directory, f) for f in sorted_files]

# 使用示例
directory = '.'  # 当前目录
output_file = 'merged_unicode_flash_mob.pptx'

files_to_merge = get_sorted_pptx_files(directory)
if files_to_merge:
    print(f"找到以下文件并按顺序合并：")
    for file in files_to_merge:
        print(f"  - {os.path.basename(file)}")
    
    print("\n开始合并...")
    total_slides = merge_pptx(files_to_merge, output_file)
    print(f"\n合并完成，输出文件：{output_file}")
    print(f"总共合并了 {total_slides} 张幻灯片")
else:
    print("没有找到符合条件的文件")
